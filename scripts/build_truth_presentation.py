"""Build PPTX for Canva import — Truth Behind the Image (4 ESO)."""
from io import BytesIO
from pathlib import Path

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(r"c:\Users\mathi\Downloads\Truth-Behind-the-Image-Canva.pptx")

IMAGES = {
    "title": "https://images.unsplash.com/photo-1427504494785-3a9ca7044f45?w=1200&q=80",
    "story": "https://images.unsplash.com/photo-1456513080510-7bf3a84b82f8?w=1200&q=80",
    "where": "https://images.unsplash.com/photo-1524661135-423995f22d0b?w=1200&q=80",
    "impact": "https://images.unsplash.com/photo-1509062522246-3755977927d7?w=1200&q=80",
    "voices": "https://images.unsplash.com/photo-1504711434969-e33886168f5c?w=1200&q=80",
    "sources": "https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=1200&q=80",
    "reflection": "https://images.unsplash.com/photo-1522202176988-66273c2fd55f?w=1200&q=80",
}

BLUE = RGBColor(0, 51, 102)
GRAY = RGBColor(80, 80, 80)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(245, 248, 252)


def fetch_image(url: str) -> BytesIO:
    r = requests.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
    r.raise_for_status()
    return BytesIO(r.content)


def add_footer(slide, prs, text: str):
    box = slide.shapes.add_textbox(Inches(0.4), prs.slide_height - Inches(0.45), Inches(9), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY


def add_title_block(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(6.2), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(6)


def add_bullets(slide, lines: list[str], top=1.35):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(6.3), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(8)


def add_image_right(slide, key: str, cache: dict):
    if key not in cache:
        cache[key] = fetch_image(IMAGES[key])
    slide.shapes.add_picture(cache[key], Inches(6.95), Inches(0.35), width=Inches(6.0))


def slide_base(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    # accent bar left
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    return slide


def main():
    cache: dict[str, BytesIO] = {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    s = slide_base(prs)
    add_image_right(s, "title", cache)
    add_title_block(
        s,
        "One image,\nthousands of questions",
        "School and war — what social media does not show",
    )
    box = s.shapes.add_textbox(Inches(0.55), Inches(2.5), Inches(6), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Pau & Martina"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.text = "Institut La Salle Reus · 4º ESO · English"
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p3 = tf.add_paragraph()
    p3.text = "Truth Behind the Image · 2025–26"
    p3.font.size = Pt(12)
    p3.font.color.rgb = GRAY
    add_footer(s, prs, "Slide 1 — Title")

    # 2 Story
    s = slide_base(prs)
    add_image_right(s, "story", cache)
    add_title_block(s, "What is this about?")
    add_bullets(
        s,
        [
            "We see many war photos on TikTok and Instagram. Some are real. Some are not.",
            "In February 2022 fighting started again in Ukraine. It is still going on.",
            "Millions of civilians left their homes. Many are children.",
            "Some posts use old photos or propaganda. Our project is about education.",
        ],
    )
    add_footer(s, prs, "Slide 2")

    # 3 Where
    s = slide_base(prs)
    add_image_right(s, "where", cache)
    add_title_block(s, "Where and when?")
    add_bullets(
        s,
        [
            "Place: Ukraine (Eastern Europe). Poland and Romania also received refugees.",
            "When: From 24 February 2022 until today. Short ceasefire talks, no lasting peace.",
            "Context: A big humanitarian crisis recognised by the UN and NGOs.",
        ],
    )
    add_footer(s, prs, "Slide 3")

    # 4 Impact
    s = slide_base(prs)
    add_image_right(s, "impact", cache)
    add_title_block(s, "Human impact")
    add_bullets(
        s,
        [
            "Displacement: families moved. Children changed school many times.",
            "Shortages: less food, water and electricity. Schools closed when unsafe.",
            "Trauma: loud attacks and long time in shelters are hard for kids.",
            "Resilience: teachers used online classes. NGOs opened learning centres.",
        ],
    )
    add_footer(s, prs, "Slide 4")

    # 5 Voices
    s = slide_base(prs)
    add_image_right(s, "voices", cache)
    add_title_block(s, "Voices and vocabulary")
    quote = s.shapes.add_shape(1, Inches(0.55), Inches(1.25), Inches(6.2), Inches(1.35))
    quote.fill.solid()
    quote.fill.fore_color.rgb = LIGHT
    quote.line.color.rgb = BLUE
    tf = quote.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = '"You cannot teach a child who is too afraid to enter the classroom."'
    p.font.italic = True
    p.font.size = Pt(13)
    p2 = tf.add_paragraph()
    p2.text = "— UNICEF (2023)"
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    box = s.shapes.add_textbox(Inches(0.55), Inches(2.75), Inches(6.3), Inches(1.2))
    tf2 = box.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Reported speech: The BBC said that some viral photos from Ukraine were taken years earlier in Syria or even in video games."
    p.font.size = Pt(13)
    box2 = s.shapes.add_textbox(Inches(0.55), Inches(4.0), Inches(6.3), Inches(2.5))
    tf3 = box2.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "Vocabulary: displacement · ceasefire · humanitarian crisis · civilians · resilience · shortages · trauma · propaganda"
    p.font.size = Pt(12)
    p.font.color.rgb = BLUE
    add_footer(s, prs, "Slide 5")

    # 6 Sources
    s = slide_base(prs)
    add_image_right(s, "sources", cache)
    add_title_block(s, "How we checked our sources")
    add_bullets(
        s,
        [
            "We read BBC, Reuters and UNICEF — not random blogs.",
            "Same story in two or three sites? If not, we were careful.",
            "We checked the date of the photo. Before 2022 = not from this war.",
            "We used calm photos only (no violent images).",
        ],
        top=1.2,
    )
    box = s.shapes.add_textbox(Inches(0.55), Inches(5.5), Inches(6.5), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sources: BBC News, Reuters, UNICEF Ukraine education reports (2023–24)."
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = GRAY
    add_footer(s, prs, "Slide 6")

    # 7 Reflection
    s = slide_base(prs)
    add_image_right(s, "reflection", cache)
    add_title_block(s, "Our reflection")
    add_bullets(
        s,
        [
            "Why is misinformation dangerous?",
            "→ People donate to the wrong place or stop helping.",
            "→ Fake news creates more hate between groups.",
            "What can young people do?",
            "→ Pause before you share. Who published this?",
            "→ Follow UNICEF / Red Cross. Explain with facts, not rumours.",
            "→ In La Salle we can do a short talk or poster at school.",
            "Thank you for listening!",
        ],
        top=1.15,
    )
    add_footer(s, prs, "Slide 7")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
